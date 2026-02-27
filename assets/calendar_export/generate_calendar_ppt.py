#!/usr/bin/env python3
"""
Générateur de Planning de Garde — Export PPT
Génère un calendrier mensuel illustré à partir d'un CSV de garde.

Usage:
    python3 generate_calendar_ppt.py --csv planning.csv --output output.pptx --assets ./
    python3 generate_calendar_ppt.py --csv planning.csv  # output = même dossier que le CSV
"""

import argparse
import calendar
import os
import sys
from datetime import timedelta

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# ── Configuration ──────────────────────────────────────────────────────

PAPA_KEY = "robert"
MAMAN_KEY = "justine"

MONTH_NAMES = [
    "", "Janvier", "Février", "Mars", "Avril", "Mai", "Juin",
    "Juillet", "Août", "Septembre", "Octobre", "Novembre", "Décembre",
]

DAY_NAMES = ["Lundi", "Mardi", "Mercredi", "Jeudi", "Vendredi", "Samedi", "Dimanche"]

BIRTHDAYS = {
    (1, 15): "Marcel",
    (3, 17): "Victoire",
    (2, 15): "Mamie Coralie",
    (9, 23): "Tatie Chloé",
    (11, 16): "Papa",
    (12, 25): "Mémé",
}

# Dégradés gauche→droite pour chaque jour de la semaine
# (couleur gauche, couleur droite) — extraits du PDF de référence Canva
LOUP_GRADIENT_LEFT = [
    RGBColor(52, 206, 52),     # Lundi     — vert foncé
    RGBColor(255, 72, 4),      # Mardi     — orange-rouge
    RGBColor(255, 107, 180),   # Mercredi  — rose vif
    RGBColor(33, 146, 255),    # Jeudi     — bleu vif
    RGBColor(255, 143, 0),     # Vendredi  — ambre foncé
    RGBColor(141, 71, 20),     # Samedi    — brun foncé
    RGBColor(251, 0, 7),       # Dimanche  — rouge vif
]

LOUP_GRADIENT_RIGHT = [
    RGBColor(151, 254, 151),   # Lundi     — vert pastel
    RGBColor(255, 158, 120),   # Mardi     — pêche
    RGBColor(255, 189, 202),   # Mercredi  — rose pastel
    RGBColor(133, 205, 250),   # Jeudi     — bleu pastel
    RGBColor(255, 212, 0),     # Vendredi  — jaune doré
    RGBColor(203, 131, 62),    # Samedi    — brun clair
    RGBColor(152, 0, 203),     # Dimanche  — violet
]

# Couleurs des en-têtes de jours (texte)
LOUP_HEADER_COLORS = [
    RGBColor(50, 205, 50),     # Lundi     — lime green
    RGBColor(255, 69, 0),      # Mardi     — orange-red
    RGBColor(255, 105, 180),   # Mercredi  — hot pink
    RGBColor(30, 144, 255),    # Jeudi     — dodger blue
    RGBColor(255, 140, 0),     # Vendredi  — dark orange
    RGBColor(139, 69, 19),     # Samedi    — saddle brown
    RGBColor(255, 0, 0),       # Dimanche  — rouge
]

MONTH_TITLE_COLOR = RGBColor(0, 151, 178)   # Teal — couleur du titre du mois
NIGHT_BAND_COLOR = RGBColor(10, 10, 80)      # Bleu nuit très foncé


# ── Données ────────────────────────────────────────────────────────────

def load_data(csv_path: str, start_from: str = "2026-03-01") -> tuple:
    """Charge le CSV de garde et retourne (calendar_data dict, DataFrame filtré)."""
    for sep in [",", ";", "\t"]:
        try:
            df = pd.read_csv(csv_path, sep=sep, encoding="utf-8-sig", parse_dates=["Date début"], dayfirst=True)
            if "Date début" in df.columns:
                break
        except Exception:
            continue
    else:
        raise ValueError(f"Impossible de lire le CSV : {csv_path}")

    df = df.dropna(subset=["Date début"])
    if start_from:
        df = df[df["Date début"] >= start_from]
    df["DateStr"] = df["Date début"].dt.strftime("%Y-%m-%d")

    cal_data = {}
    for _, row in df.iterrows():
        date_key = row["DateStr"]
        parent = str(row["Parent"]).strip() if pd.notna(row.get("Parent")) else ""
        garde_type = str(row["Type de garde"]).strip() if pd.notna(row.get("Type de garde")) else ""
        cal_data[date_key] = {"parent": parent, "type": garde_type}

    return cal_data, df


def get_parent_for_night(current_date, info_today: dict, calendar_data: dict) -> str:
    """Détermine le parent pour la NUIT.

    Le CSV indique le parent de journée, mais le calendrier affiche
    le parent chez qui les enfants dorment. Règles hors vacances :
      - Mardi soir   → Robert  (car mercredi journée chez Papa)
      - Mercredi soir → Justine (même si journée chez Papa)
      - Dimanche soir → Justine (toujours)
      - Autres soirs  → suit le parent du CSV
    En vacances, le parent CSV est directement le parent de nuit.
    """
    parent_today = info_today.get("parent", "")
    type_today = info_today.get("type", "")

    # Vacances → le parent du CSV a les enfants jour ET nuit
    if type_today == "Vacances":
        return parent_today

    weekday = current_date.weekday()  # 0=Lundi … 6=Dimanche

    # Dimanche soir → toujours chez Maman
    if weekday == 6:
        return "Justine"

    # Mercredi soir → chez Maman (même si journée chez Papa)
    if weekday == 2:
        return "Justine"

    # Mardi soir → chez Papa si mercredi il a les enfants en journée
    if weekday == 1:
        next_day = current_date + timedelta(days=1)
        next_str = next_day.strftime("%Y-%m-%d")
        info_next = calendar_data.get(next_str, {})
        if "Robert" in info_next.get("parent", ""):
            return "Robert"
        return parent_today

    return parent_today


# ── Styles PPTX ────────────────────────────────────────────────────────

def apply_gradient_fill(shape, left_color: RGBColor, right_color: RGBColor):
    """Remplit une forme avec un dégradé gauche→droite."""
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = 0  # 0° = gauche vers droite
    gs = fill.gradient_stops
    gs[0].position = 0.0
    gs[0].color.rgb = left_color
    gs[1].position = 1.0
    gs[1].color.rgb = right_color


def apply_solid_fill(shape, color: RGBColor):
    """Remplit une forme avec une couleur unie."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_image_safe(slide, img_path, left, top, height):
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, left, top, height=height)


# ── Génération des slides ──────────────────────────────────────────────

def create_month_slide(prs, year: int, month: int, calendar_data: dict, assets_dir: str):
    """Crée un slide mensuel dans la présentation."""
    img_papa = os.path.join(assets_dir, "papa.png")
    img_maman = os.path.join(assets_dir, "maman.png")
    img_ecole = os.path.join(assets_dir, "logo_ecole.png")
    img_vacances = os.path.join(assets_dir, "logo_vacances.png")
    img_anniv = os.path.join(assets_dir, "logo_anniv.png")
    img_lune = os.path.join(assets_dir, "logo_lune.png")
    img_we = os.path.join(assets_dir, "logo_we.png")

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ── Marges d'impression (≈ 0.35in de chaque côté) ──
    margin_x = Inches(0.35)
    margin_bottom = Inches(0.25)
    usable_w = prs.slide_width - 2 * margin_x
    col_w = usable_w // 7          # ≈ 1.57in par colonne
    start_x = margin_x

    # ── Titre du mois (pas de fond, texte teal) ──
    t_box = slide.shapes.add_textbox(
        start_x, Inches(0.15), usable_w, Inches(0.7),
    )
    t_box.text_frame.text = f"{MONTH_NAMES[month]} {year}"
    p = t_box.text_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = MONTH_TITLE_COLOR

    # ── En-tête jours ──
    header_y = Inches(0.9)

    for i, day_name in enumerate(DAY_NAMES):
        box = slide.shapes.add_textbox(
            start_x + i * col_w, header_y, col_w, Inches(0.3),
        )
        p = box.text_frame.paragraphs[0]
        p.text = day_name
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = LOUP_HEADER_COLORS[i]

    # ── Grille ──
    cal = calendar.monthcalendar(year, month)
    num_rows = len(cal)    # 4, 5 ou 6 semaines
    grid_top = Inches(1.25)
    grid_bottom = prs.slide_height - margin_bottom
    available_h = grid_bottom - grid_top
    row_h = available_h // num_rows   # s'adapte au nombre de semaines
    cell_gap = Inches(0.06)            # espace entre cases

    for row_idx, week in enumerate(cal):
        for col_idx, day in enumerate(week):
            if day == 0:
                continue

            current_date = pd.Timestamp(year=year, month=month, day=day)
            date_str = current_date.strftime("%Y-%m-%d")
            info = calendar_data.get(date_str, {})

            gtype = info.get("type", "")
            is_vacation = gtype == "Vacances"

            dodo_text = get_parent_for_night(current_date, info, calendar_data)

            left = start_x + col_idx * col_w
            top = grid_top + row_idx * row_h
            cell_w = col_w - cell_gap
            cell_h = row_h - cell_gap

            # Case jour (rectangle arrondi, dégradé gauche→droite)
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, cell_w, cell_h,
            )
            shape.adjustments[0] = 0.12
            apply_gradient_fill(
                shape,
                LOUP_GRADIENT_LEFT[col_idx],
                LOUP_GRADIENT_RIGHT[col_idx],
            )
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.width = Pt(2)

            # Numéro du jour
            tb_n = slide.shapes.add_textbox(
                left + Inches(0.05), top, Inches(0.5), Inches(0.35),
            )
            p_n = tb_n.text_frame.paragraphs[0]
            p_n.text = str(day)
            p_n.font.size = Pt(20)
            p_n.font.bold = True
            p_n.font.color.rgb = RGBColor(255, 255, 255)

            # ── Bandeau nuit ──
            band_h = int(cell_h * 0.35)
            band_top = top + cell_h - band_h

            band = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, band_top, cell_w, band_h,
            )
            band.adjustments[0] = 0.12
            apply_solid_fill(band, NIGHT_BAND_COLOR)
            band.line.fill.background()

            # Lune
            moon_size = int(band_h * 1.2)
            moon_top = band_top + band_h - moon_size + Inches(0.03)
            moon_left = left + Inches(0.08)

            if os.path.exists(img_lune):
                add_image_safe(slide, img_lune, moon_left, moon_top, moon_size)
            else:
                tb_m = slide.shapes.add_textbox(left, band_top, Inches(0.4), band_h)
                tb_m.text_frame.text = "🌙"

            # ── Avatars parents (nuit) ──
            dodo_lower = dodo_text.lower()
            has_papa = PAPA_KEY in dodo_lower or "papa" in dodo_lower
            has_maman = MAMAN_KEY in dodo_lower or "maman" in dodo_lower

            avatar_size = int(band_h * 1.2)
            base_right_pos = left + cell_w - avatar_size - Inches(0.04)
            avatar_top = band_top + band_h - avatar_size + Inches(0.03)

            if has_papa and has_maman:
                add_image_safe(slide, img_maman, base_right_pos, avatar_top, avatar_size)
                add_image_safe(slide, img_papa, base_right_pos - avatar_size + Inches(0.08), avatar_top, avatar_size)
            elif has_papa:
                add_image_safe(slide, img_papa, base_right_pos, avatar_top, avatar_size)
            elif has_maman:
                add_image_safe(slide, img_maman, base_right_pos, avatar_top, avatar_size)
            else:
                tb_txt = slide.shapes.add_textbox(
                    left + Inches(0.35), band_top, Inches(0.9), band_h,
                )
                p_txt = tb_txt.text_frame.paragraphs[0]
                p_txt.text = dodo_text
                p_txt.font.color.rgb = RGBColor(255, 255, 255)
                p_txt.font.size = Pt(7)
                p_txt.alignment = PP_ALIGN.RIGHT

            # ── Icône activité ──
            icon_file = None
            if BIRTHDAYS.get((month, day)):
                icon_file = img_anniv
            elif is_vacation:
                icon_file = img_vacances
            elif col_idx in [2, 5, 6]:  # Mercredi, Samedi, Dimanche
                icon_file = img_we
            elif col_idx in [0, 1, 3, 4]:  # Jours d'école
                icon_file = img_ecole

            if icon_file and os.path.exists(icon_file):
                add_image_safe(
                    slide, icon_file,
                    left + cell_w - Inches(0.5), top + Inches(0.03),
                    Inches(0.42),
                )
            elif icon_file:
                tb_f = slide.shapes.add_textbox(
                    left + cell_w - Inches(0.5), top, Inches(0.5), Inches(0.5),
                )
                if "anniv" in icon_file:
                    tb_f.text_frame.text = "🎂"
                elif "we" in icon_file:
                    tb_f.text_frame.text = "🎈"
                elif "vacances" in icon_file:
                    tb_f.text_frame.text = "🪁"
                else:
                    tb_f.text_frame.text = "🎒"


# ── Main ───────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Génère un planning de garde PPT")
    parser.add_argument("--csv", required=True, help="Chemin du fichier CSV de garde")
    parser.add_argument("--output", help="Chemin du fichier PPTX de sortie (défaut : même dossier que le CSV)")
    parser.add_argument("--assets", help="Dossier contenant les images (papa.png, maman.png, etc.)")
    parser.add_argument("--start-from", default="2026-03-01", help="Date de début (YYYY-MM-DD), défaut : 2026-03-01")
    args = parser.parse_args()

    csv_path = os.path.abspath(args.csv)
    if not os.path.exists(csv_path):
        print(f"❌ CSV introuvable : {csv_path}", file=sys.stderr)
        sys.exit(1)

    # Assets : même dossier que ce script par défaut
    assets_dir = args.assets or os.path.dirname(os.path.abspath(__file__))

    # Output : même dossier que le CSV par défaut
    if args.output:
        output_path = os.path.abspath(args.output)
    else:
        csv_dir = os.path.dirname(csv_path)
        output_path = os.path.join(csv_dir, "Calendrier_Garde.pptx")

    print(f"📅 Chargement CSV : {csv_path}")
    calendar_data, df = load_data(csv_path, start_from=args.start_from)

    print(f"📊 {len(calendar_data)} jours de données chargés")

    # Créer la présentation (paysage A4)
    prs = Presentation()
    prs.slide_width = Inches(11.69)
    prs.slide_height = Inches(8.27)

    # Générer les mois
    months = sorted(set((d.year, d.month) for d in pd.to_datetime(list(calendar_data.keys()))))
    if not months:
        print("⚠️ Aucun mois à générer (CSV vide ou dates filtrées)")
        sys.exit(0)

    for y, m in months:
        print(f"  → {MONTH_NAMES[m]} {y}")
        create_month_slide(prs, y, m, calendar_data, assets_dir)

    prs.save(output_path)
    # Afficher le chemin sur stdout (lu par l'app Flutter)
    print(f"✅ {output_path}")


if __name__ == "__main__":
    main()
